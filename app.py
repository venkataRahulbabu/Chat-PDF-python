import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_community.vectorstores.faiss import FAISS
from dotenv import load_dotenv

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Helper functions
def extract_text_from_pdf(file):
    """Extract text from a PDF file."""
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file):
    """Extract text from a Word (.docx) file."""
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file):
    """Extract text from a PowerPoint (.pptx) file."""
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text

def extract_text_from_txt(file):
    """Extract text from a plain text (.txt) file."""
    return file.read().decode("utf-8")

def extract_text_from_file(file):
    """Extract text based on file type."""
    if file.name.endswith(".pdf"):
        return extract_text_from_pdf(file)
    elif file.name.endswith(".docx"):
        return extract_text_from_docx(file)
    elif file.name.endswith(".pptx"):
        return extract_text_from_pptx(file)
    elif file.name.endswith(".txt"):
        return extract_text_from_txt(file)
    else:
        return ""

def get_text_chunks(text):
    """Split text into manageable chunks."""
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    """Create vector store from text chunks."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    """Set up a conversational chain for QA."""
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the question is not based on the context or is a general comment, handle it appropriately.\n\n
    Context:\n{context}\n
    Question:\n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

def is_general_comment(user_question):
    """Identify general comments."""
    general_comments = ["well done", "thank you", "thanks", "good job", "great", "nice", "awesome"]
    return any(comment in user_question.lower() for comment in general_comments)

def user_input(user_question):
    """Process user input and return response."""
    if is_general_comment(user_question):
        return "You're welcome! I'm here to help you with any questions you have."
    
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()

    response = chain.invoke({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

# Main Streamlit app
def main():
    st.set_page_config("Chat Files", layout="wide")
    st.header("Chat with Files")

    if "conversation" not in st.session_state:
        st.session_state.conversation = []

    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader(
            "Upload your Files (PDFs, Word Docs, PowerPoint, Text) and Click on Submit & Process", 
            accept_multiple_files=True,
            type=["pdf", "docx", "pptx", "txt"]
        )
        if st.button("Submit & Process"):
            if uploaded_files:
                with st.spinner("Processing Files..."):
                    raw_text = ""
                    for file in uploaded_files:
                        raw_text += extract_text_from_file(file)
                    
                    if raw_text.strip():
                        text_chunks = get_text_chunks(raw_text)
                        get_vector_store(text_chunks)
                        st.success("Processing complete. You can now ask questions!")
                    else:
                        st.error("No text could be extracted from the uploaded files.")
            else:
                st.error("Please upload at least one file.")

    # Custom CSS for chat styling
    st.markdown("""
        <style>
            .chat-container {
                display: flex;
                flex-direction: column;
            }
            .user-message, .model-response {
                padding: 10px;
                border-radius: 10px;
                margin-bottom: 10px;
                max-width: 60%;
            }
            .user-message {
                align-self: flex-end;
                background-color: #DCF8C6;
                font-weight: 600;
            }
            .model-response {
                align-self: flex-start;
                background-color: #F1F0F0;
                font-weight: normal;
            }
        </style>
    """, unsafe_allow_html=True)

    # Display conversation
    for i, (question, answer) in enumerate(st.session_state.conversation):
        st.markdown(f'<div class="chat-container"><div class="user-message">YOU: {question}</div><div class="model-response">MODEL: {answer}</div></div>', unsafe_allow_html=True)

    new_question = st.text_input("Your new question:", key=f"question_{len(st.session_state.conversation)}")

    if st.button("Submit Query"):
        if new_question:
            answer = user_input(new_question)
            st.session_state.conversation.append((new_question, answer))
            st.rerun()

if __name__ == "__main__":
    main()
